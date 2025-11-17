"""RAG Feasibility Study Generator
-----------------------------------

This standalone FastAPI service ingests up to 1 GB of project materials,
builds a FAISS vector index, and generates a grounded feasibility study
using an LLM. The implementation follows the blueprint previously used in
the Streamlit AI workspace so teams can restore the RAG workflow outside the
app.

Key endpoints:
    POST /collect  -> store financial_snapshot/cell_map from the Excel model
    POST /ingest   -> streamed file uploads (PDF/DOCX/PPTX/TXT/CSV/XLSX)
    POST /generate -> draft feasibility sections with citations

To run locally:
    export OPENAI_API_KEY="sk-..."
    python rag_app.py
"""

from __future__ import annotations

import os
import io
import json
import hashlib
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional

from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.responses import JSONResponse
from pydantic import BaseModel

# Third-party libraries (install via requirements in the README snippet)
import pandas as pd

from docx import Document as DocxDocument
from pptx import Presentation
import pypdf

import faiss
from sentence_transformers import SentenceTransformer

try:  # Optional reranker
    from sentence_transformers import CrossEncoder

    _HAS_RERANK = True
except Exception:  # pragma: no cover - fallback when package absent
    CrossEncoder = None
    _HAS_RERANK = False

# OpenAI client (Chat Completions)
try:
    from openai import OpenAI

    _HAS_OPENAI = True
except Exception:  # pragma: no cover
    OpenAI = None
    _HAS_OPENAI = False

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter

# -------------------- Config --------------------
EMBED_MODEL_NAME = os.getenv("EMBED_MODEL", "sentence-transformers/all-MiniLM-L6-v2")
RERANK_MODEL_NAME = os.getenv("RERANK_MODEL", "cross-encoder/ms-marco-MiniLM-L-6-v2")

CHUNK_TOKENS = int(os.getenv("CHUNK_TOKENS", "500"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "100"))
TOP_K = int(os.getenv("TOP_K", "12"))
RERANK_K = int(os.getenv("RERANK_K", "5"))

DATA_DIR = os.getenv("DATA_DIR", "./projects")

LLM_MODEL_DEFAULT = os.getenv("LLM_MODEL", "gpt-4o")
TEMPERATURE = float(os.getenv("LLM_TEMPERATURE", "0.2"))

DEFAULT_SECTIONS = [
    "Executive Summary",
    "Project Description & Scope",
    "Market & Demand Analysis",
    "Technical & Operations",
    "Legal, Permitting & Environmental",
    "Implementation Plan",
    "Financial Analysis",
    "Risk Assessment & ESG",
    "Conclusion & Recommendation",
    "Appendices",
]

SYSTEM_PROMPT = (
    "You are a financial analyst producing a feasibility study. Only use facts from provided "
    "CONTEXT and FINANCIAL_SNAPSHOT. Cite sources inline like [Source: filename p.12] or "
    "[Sheet: Assumptions!B7]. If a claim is unsupported, say so. Keep each section structured, "
    "concise, and decision-oriented."
)

SECTION_PROMPT = """[GOAL]
Draft the {section} for the feasibility study.

[GUIDANCE]
- Use FINANCIAL_SNAPSHOT metrics explicitly where relevant.
- Use CONTEXT passages with inline citations [Source: <file> p.<n>] or [Sheet: <name>!<cell>].
- State uncertainties and missing data.
- Avoid boilerplate; keep it specific to the project.

[FINANCIAL_SNAPSHOT]
{fin}

[CONTEXT]
{ctx}

[OUTPUT]
- 3–7 well-structured paragraphs
- Subheadings
- Bullet lists for key metrics & risks
- Citations inline
"""

# -------------------- Utilities --------------------


def sha256_file(path: str) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def ensure_dirs(project_id: str) -> str:
    base = os.path.join(DATA_DIR, project_id)
    os.makedirs(os.path.join(base, "uploads"), exist_ok=True)
    os.makedirs(os.path.join(base, "parsed"), exist_ok=True)
    os.makedirs(os.path.join(base, "index"), exist_ok=True)
    os.makedirs(os.path.join(base, "financial"), exist_ok=True)
    os.makedirs(os.path.join(base, "charts"), exist_ok=True)
    return base


async def stream_save(upload: UploadFile, dest_path: str):
    with open(dest_path, "wb") as out:
        while True:
            chunk = await upload.read(1024 * 1024)
            if not chunk:
                break
            out.write(chunk)


def simple_tokenizer(text: str) -> List[str]:
    return text.split()


def chunk_text(text: str, chunk_tokens=CHUNK_TOKENS, overlap=CHUNK_OVERLAP) -> List[str]:
    tokens = simple_tokenizer(text)
    chunks: List[str] = []
    start = 0
    while start < len(tokens):
        end = min(len(tokens), start + chunk_tokens)
        chunk = " ".join(tokens[start:end])
        chunks.append(chunk)
        if end == len(tokens):
            break
        start = max(0, end - overlap)
    return chunks


# -------------------- Parsing --------------------


def parse_pdf(path: str) -> List[Dict[str, Any]]:
    reader = pypdf.PdfReader(path)
    items: List[Dict[str, Any]] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        items.append({"page_or_sheet": i, "text": text})
    return items


def parse_docx(path: str) -> List[Dict[str, Any]]:
    doc = DocxDocument(path)
    text = "\n".join([p.text for p in doc.paragraphs])
    return [{"page_or_sheet": 1, "text": text}]


def parse_pptx(path: str) -> List[Dict[str, Any]]:
    prs = Presentation(path)
    items: List[Dict[str, Any]] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        items.append({"page_or_sheet": i, "text": "\n".join(texts)})
    return items


def parse_txt(path: str) -> List[Dict[str, Any]]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return [{"page_or_sheet": 1, "text": f.read()}]


# -------------------- Vector Index --------------------


class VectorIndex:
    def __init__(self, project_dir: str):
        self.project_dir = project_dir
        self.index_dir = os.path.join(project_dir, "index")
        self.index_path = os.path.join(self.index_dir, "faiss.index")
        self.meta_path = os.path.join(self.index_dir, "meta.jsonl")
        self.model = SentenceTransformer(EMBED_MODEL_NAME)
        self.dim = self.model.get_sentence_embedding_dimension()
        self.index = None
        self._load()

    def _load(self):
        if os.path.exists(self.index_path):
            self.index = faiss.read_index(self.index_path)
        else:
            self.index = faiss.IndexFlatIP(self.dim)
        if not os.path.exists(self.meta_path):
            open(self.meta_path, "a").close()

    def _save(self):
        faiss.write_index(self.index, self.index_path)

    def add(self, texts: List[str], metas: List[Dict[str, Any]]):
        if not texts:
            return
        embeddings = self.model.encode(texts, normalize_embeddings=True)
        self.index.add(embeddings)
        with open(self.meta_path, "a", encoding="utf-8") as f:
            for m, t in zip(metas, texts):
                f.write(json.dumps({"meta": m, "text": t}) + "\n")
        self._save()

    def search(self, query: str, top_k: int = TOP_K) -> List[Dict[str, Any]]:
        if self.index.ntotal == 0:
            return []
        q = self.model.encode([query], normalize_embeddings=True)
        sims, idxs = self.index.search(q, top_k)
        results: List[Dict[str, Any]] = []
        with open(self.meta_path, "r", encoding="utf-8") as f:
            lines = f.readlines()
        for i in idxs[0]:
            if i < 0 or i >= len(lines):
                continue
            rec = json.loads(lines[i])
            results.append(rec)
        return results


class Reranker:
    def __init__(self, name=RERANK_MODEL_NAME):
        self.enabled = False
        if _HAS_RERANK:
            try:
                self.model = CrossEncoder(name)
                self.enabled = True
            except Exception:
                self.model = None
        else:
            self.model = None

    def rerank(self, query: str, passages: List[Dict[str, Any]], k=RERANK_K):
        if not self.enabled or not passages:
            return passages[:k]
        pairs = [(query, p["text"]) for p in passages]
        scores = self.model.predict(pairs)
        ranked = sorted(zip(passages, scores), key=lambda x: x[1], reverse=True)
        return [p for p, _ in ranked[:k]]


class RagLLMClient:
    """LLM wrapper that prefers OpenAI SDK when present and falls back to REST."""

    def __init__(self):
        self.model = LLM_MODEL_DEFAULT
        self.api_key = os.getenv("OPENAI_API_KEY")
        self.has_sdk = _HAS_OPENAI and self.api_key
        if self.has_sdk:
            self.client = OpenAI(api_key=self.api_key)
        else:
            self.client = None

    def complete(self, prompt: str, system: Optional[str] = None) -> str:
        if self.client is None:
            raise RuntimeError("OpenAI client unavailable: set OPENAI_API_KEY or install openai")
        messages = []
        if system:
            messages.append({"role": "system", "content": system})
        messages.append({"role": "user", "content": prompt})
        resp = self.client.chat.completions.create(
            model=self.model, messages=messages, temperature=TEMPERATURE
        )
        return resp.choices[0].message.content


# -------------------- Charts --------------------


def _ensure_chart_dir(project_dir: str) -> str:
    chart_dir = os.path.join(project_dir, "charts")
    os.makedirs(chart_dir, exist_ok=True)
    return chart_dir


def _fmt_currency(curr: str):
    def _fmt(x, pos):  # pragma: no cover - matplotlib formatter
        try:
            return f"{curr} {x:,.0f}" if curr else f"{x:,.0f}"
        except Exception:
            return f"{x:,.0f}"

    return FuncFormatter(_fmt)


def plot_npv_curve(financial: dict, out_path: str):
    xs, ys = [], []
    if financial.get("scenarios"):
        xs = [s.get("name", f"S{i+1}") for i, s in enumerate(financial["scenarios"])]
        ys = [float(s.get("npv", 0) or 0) for s in financial["scenarios"]]
    elif financial.get("sensitivities"):
        sens = [s for s in financial["sensitivities"] if s.get("npv") is not None]
        sens = sorted(sens, key=lambda s: s.get("delta", 0))
        xs = [f"{int((s.get('delta') or 0)*100)}%" for s in sens]
        ys = [float(s.get("npv", 0) or 0) for s in sens]
    else:
        return None

    fig = plt.figure()
    plt.plot(xs, ys, marker="o")
    plt.title("NPV Curve")
    plt.xlabel("Scenario / Delta")
    plt.ylabel("NPV")
    curr = financial.get("currency", "") or ""
    plt.gca().yaxis.set_major_formatter(_fmt_currency(curr))
    plt.grid(True, linestyle="--", alpha=0.4)
    fig.tight_layout()
    fig.savefig(out_path, dpi=200)
    plt.close(fig)
    return out_path


# -------------------- API models --------------------


class IngestResponse(BaseModel):
    project_id: str
    files: List[str]


class CollectRequest(BaseModel):
    project_id: str
    financial_snapshot: Dict[str, Any]
    cell_map: Optional[Dict[str, Any]] = None
    workbook_hash: Optional[str] = None
    model_version: Optional[str] = None


class GenerateRequest(BaseModel):
    project_id: str
    section_outline: Optional[List[str]] = None
    query_hint: Optional[str] = None


# -------------------- FastAPI App --------------------

app = FastAPI(title="RAG Feasibility Study Generator")
llm_client: Optional[RagLLMClient] = None


def _llm() -> RagLLMClient:
    global llm_client
    if llm_client is None:
        llm_client = RagLLMClient()
    return llm_client


@app.post("/collect")
def collect(req: CollectRequest):
    base = ensure_dirs(req.project_id)
    snap_path = os.path.join(base, "financial", "snapshot.json")
    snapshot = dict(req.financial_snapshot)
    snapshot["collected_at"] = datetime.now(timezone.utc).isoformat()
    if req.workbook_hash:
        snapshot["workbook_hash"] = req.workbook_hash
    if req.model_version:
        snapshot["model_version"] = req.model_version

    with open(snap_path, "w", encoding="utf-8") as f:
        json.dump(snapshot, f, indent=2)

    if req.cell_map:
        with open(os.path.join(base, "financial", "cell_map.json"), "w", encoding="utf-8") as f:
            json.dump(req.cell_map, f, indent=2)

    return {"status": "ok", "project_id": req.project_id, "snapshot_path": snap_path}


@app.post("/ingest", response_model=IngestResponse)
async def ingest(project_id: str = Form(...), files: List[UploadFile] = File(...)):
    base = ensure_dirs(project_id)
    saved: List[str] = []
    for f in files:
        dest = os.path.join(base, "uploads", f.filename)
        await stream_save(f, dest)
        saved.append(dest)

        parsed_items: List[Dict[str, Any]] = []
        ext = os.path.splitext(f.filename)[1].lower()
        if ext == ".pdf":
            parsed_items = parse_pdf(dest)
        elif ext == ".docx":
            parsed_items = parse_docx(dest)
        elif ext == ".pptx":
            parsed_items = parse_pptx(dest)
        elif ext in (".txt", ".md"):
            parsed_items = parse_txt(dest)
        elif ext in (".xlsx", ".xls", ".csv"):
            try:
                if ext == ".csv":
                    df = pd.read_csv(dest)
                    text = df.to_csv(index=False)
                    parsed_items = [{"page_or_sheet": 1, "text": text}]
                else:
                    xls = pd.ExcelFile(dest)
                    for s in xls.sheet_names:
                        df = xls.parse(s)
                        text = df.to_csv(index=False)
                        parsed_items.append({"page_or_sheet": s, "text": text})
            except Exception:
                parsed_items = []

        parsed_path = os.path.join(base, "parsed", f"{os.path.basename(dest)}.jsonl")
        with open(parsed_path, "w", encoding="utf-8") as out:
            for item in parsed_items:
                out.write(json.dumps(item) + "\n")

        vi = VectorIndex(base)
        texts, metas = [], []
        for item in parsed_items:
            chunks = chunk_text(item.get("text", ""))
            for ch in chunks:
                texts.append(ch)
                metas.append(
                    {
                        "project_id": project_id,
                        "file_path": dest,
                        "file_type": ext[1:],
                        "page_or_sheet": item.get("page_or_sheet"),
                        "section": None,
                        "char_start": 0,
                        "char_end": len(ch),
                        "hash": sha256_file(dest),
                    }
                )
        vi.add(texts, metas)

        if ext in (".xlsx", ".xls"):
            snap_path = os.path.join(base, "financial", "snapshot.json")
            if not os.path.exists(snap_path):
                snap = parse_excel_metrics(dest)
                if snap:
                    with open(snap_path, "w", encoding="utf-8") as sf:
                        json.dump(snap, sf, indent=2)

    return IngestResponse(project_id=project_id, files=saved)


def parse_excel_metrics(path: str) -> Dict[str, Any]:
    wb_hash = sha256_file(path)
    try:
        xls = pd.ExcelFile(path)
    except Exception:
        return {}

    snapshot: Dict[str, Any] = {
        "as_of": datetime.now(timezone.utc).isoformat(),
        "workbook_path": path,
        "workbook_hash": wb_hash,
        "currency": None,
        "assumptions": {},
        "capex_total": None,
        "opex_annual": None,
        "revenue_annual": None,
        "npv": None,
        "irr": None,
        "payback_years": None,
        "dscr_min": None,
        "sensitivities": [],
        "scenarios": [],
    }

    keywords = ["NPV", "IRR", "DSCR", "Payback", "Capex", "Opex", "Revenue", "Discount", "Tax"]

    for sheet in xls.sheet_names:
        try:
            df = xls.parse(sheet)
        except Exception:
            continue
        df_str = df.astype(str)
        if snapshot["currency"] is None:
            if df_str.apply(
                lambda col: col.str.contains("USD|ZAR|EUR|\\$|R|€", case=False, regex=True, na=False)
            ).any().any():
                snapshot["currency"] = "DETECTED"

        for kw in keywords:
            matches = df_str.apply(lambda col: col.str.contains(fr"\b{kw}\b", case=False, regex=True, na=False)).any(
                axis=1
            )
            if matches.any():
                row_idx = matches.idxmax()
                row = df.iloc[row_idx]
                val = None
                for v in row:
                    try:
                        vv = float(str(v).replace(",", ""))
                        val = vv
                        break
                    except Exception:
                        continue
                kw_l = kw.lower()
                if kw_l == "npv":
                    snapshot["npv"] = val
                elif kw_l == "irr":
                    snapshot["irr"] = val
                elif kw_l == "dscr":
                    snapshot["dscr_min"] = val
                elif kw_l == "payback":
                    snapshot["payback_years"] = val
                elif kw_l == "capex":
                    snapshot["capex_total"] = val
                elif kw_l == "opex":
                    snapshot["opex_annual"] = val
                elif kw_l == "revenue":
                    snapshot["revenue_annual"] = val
                elif kw_l == "discount":
                    snapshot["assumptions"].setdefault("discount_rate", val)
                elif kw_l == "tax":
                    snapshot["assumptions"].setdefault("tax_rate", val)

    return snapshot


@app.post("/generate")
def generate(req: GenerateRequest):
    base = ensure_dirs(req.project_id)
    vi = VectorIndex(base)
    reranker = Reranker()

    snap_path = os.path.join(base, "financial", "snapshot.json")
    financial: Dict[str, Any] = {}
    if os.path.exists(snap_path):
        with open(snap_path, "r", encoding="utf-8") as f:
            financial = json.load(f)

    hint = req.query_hint or "feasibility study for project"
    passages = vi.search(hint, top_k=TOP_K)
    passages = reranker.rerank(hint, passages, k=RERANK_K)

    def fmt_passage(p: Dict[str, Any]) -> str:
        m = p.get("meta", {})
        fname = os.path.basename(m.get("file_path", "source"))
        page = m.get("page_or_sheet", "?")
        return f"[Source: {fname} p.{page}]\n" + p.get("text", "")

    ctx_block = "\n\n".join(fmt_passage(p) for p in passages)
    sections = req.section_outline or DEFAULT_SECTIONS

    fin_lines: List[str] = []
    for k, v in financial.items():
        if isinstance(v, (str, int, float)):
            fin_lines.append(f"- {k}: {v}")
    if financial.get("assumptions"):
        fin_lines.append("- assumptions:")
        for ak, av in financial["assumptions"].items():
            fin_lines.append(f"  - {ak}: {av}")
    fin_block = "\n".join(fin_lines) or "- (no financial snapshot found)"

    outputs: List[Dict[str, Any]] = []
    for sec in sections:
        prompt = SECTION_PROMPT.format(section=sec, fin=fin_block, ctx=ctx_block)
        try:
            txt = _llm().complete(prompt, system=SYSTEM_PROMPT)
        except Exception as exc:  # pragma: no cover - runtime path
            raise HTTPException(status_code=503, detail=str(exc))
        outputs.append({"section": sec, "content": txt})

    report = {
        "project_id": req.project_id,
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "sections": outputs,
        "financial_snapshot": financial,
    }

    out_json = os.path.join(base, "report.json")
    out_md = os.path.join(base, "report.md")

    with open(out_json, "w", encoding="utf-8") as jf:
        json.dump(report, jf, indent=2)

    chart_dir = _ensure_chart_dir(base)
    npv_png = os.path.join(chart_dir, "npv_curve.png")
    try:
        plot_npv_curve(financial, npv_png)
    except Exception:
        npv_png = None

    with open(out_md, "w", encoding="utf-8") as mf:
        mf.write("# Feasibility Study\n\n")
        mf.write(f"**Project:** {req.project_id}\n\n")
        mf.write(f"**Generated at:** {report['generated_at']}\n\n")
        for sec in outputs:
            mf.write(f"## {sec['section']}\n\n{sec['content']}\n\n")
        mf.write("---\n\n### Financial Snapshot\n\n")
        mf.write("```json\n" + json.dumps(financial, indent=2) + "\n```\n\n")
        if npv_png and os.path.exists(npv_png):
            mf.write(f"![NPV Curve](charts/{os.path.basename(npv_png)})\n\n")

    return JSONResponse(
        {
            "project_id": req.project_id,
            "report_json": out_json,
            "report_md": out_md,
            "sections": [s["section"] for s in outputs],
        }
    )


if __name__ == "__main__":  # pragma: no cover
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)

